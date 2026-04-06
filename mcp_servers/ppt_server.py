"""
DeckGenius AI - PowerPoint MCP Server
Module: ppt_server.py

Description:
A Model Context Protocol (MCP) server that exposes PowerPoint generation tools to the LLM.
It leverages the `python-pptx` library to programmatically assemble `.pptx` documents locally
based on the instructions provided by the Groq Llama 3 API. Architecture follows strict
Object-Oriented principles for efficient state management and concurrent requests handling.
"""

from mcp.server.fastmcp import FastMCP
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Initialize the MCP Server namespace for slides generation
mcp = FastMCP("PPTOperations")


class PPTStyle:
    """Lightweight default styling for generated decks.

    Note: python-pptx can't apply a full PowerPoint theme reliably without a template,
    so this provides consistent in-code formatting as a good default.
    """

    BG = RGBColor(248, 250, 252)          # near-slate-50
    ACCENT = RGBColor(99, 102, 241)       # indigo
    TITLE = RGBColor(15, 23, 42)          # slate-900
    BODY = RGBColor(51, 65, 85)           # slate-700
    MUTED = RGBColor(100, 116, 139)       # slate-500

    TITLE_SIZE = Pt(34)
    SUBTITLE_SIZE = Pt(18)
    BODY_SIZE = Pt(20)


def _apply_slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = PPTStyle.BG


def _add_accent_bar(slide) -> None:
    # A slim accent bar on the left gives a "designed" look without needing templates.
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.2),
        Inches(0.7),
        Inches(0.18),
        Inches(6.0),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PPTStyle.ACCENT
    bar.line.fill.background()  # no border


def _style_title_shape(title_shape) -> None:
    if not title_shape or not getattr(title_shape, "text_frame", None):
        return
    tf = title_shape.text_frame
    if not tf.paragraphs:
        return
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = PPTStyle.TITLE_SIZE
    p.font.color.rgb = PPTStyle.TITLE


def _style_body_text_frame(text_frame) -> None:
    for p in text_frame.paragraphs:
        p.font.size = PPTStyle.BODY_SIZE
        p.font.color.rgb = PPTStyle.BODY


class PPTManager:
    """
    PPTManager: A singleton-like class to manage in-memory presentation states.
    Prevents repetitive disk reads by caching the active `Presentation` object.
    Automatically resolves exact path outputs so generated files don't get lost.
    """
    def __init__(self):
        self._presentation = None

    def get_absolute_path(self, filename: str) -> str:
        """
        Calculates the absolute path to save the PPTX file consistently.
        Ensures output goes to the `generated_presentations` directory at the project root.
        """
        # Resolve base directory (parent of `mcp_servers/`)
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        output_dir = os.path.join(base_dir, "generated_presentations")
        
        # Ensure the target directory exists before saving
        os.makedirs(output_dir, exist_ok=True)
        return os.path.join(output_dir, os.path.basename(filename))

    def load_or_create(self, filename: str):
        """
        Attempts to load an existing `.pptx` file from disk; creates a new one if it fails.
        Useful when the LLM generates a presentation over sequential tool calls.
        """
        abs_path = self.get_absolute_path(filename)
        if self._presentation is None:
            try:
                # Try loading an existing file (if the MCP server rebooted mid-generation)
                self._presentation = Presentation(abs_path)
            except Exception:
                # Otherwise, initialize a fresh canvas
                self._presentation = Presentation()

    def create_from_template(self, filename: str, template_path: str):
        abs_out = self.get_absolute_path(filename)
        self._presentation = Presentation(template_path)
        self._presentation.save(abs_out)
        return abs_out

# Instantiate the singleton instance for the global server scope
ppt_manager = PPTManager()

@mcp.tool()
def create_presentation(filename: str) -> str:
    """
    MCP Tool: create_presentation
    Initializes a new empty PowerPoint file. Must be called first in any PPT workflow.
    """
    ppt_manager._presentation = Presentation()
    abs_path = ppt_manager.get_absolute_path(filename)
    try:
        ppt_manager._presentation.save(abs_path)
        return f"? SUCCESS: Created new presentation at '{abs_path}'"
    except Exception as e:
        return f"? ERROR: Failed to create presentation. Details: {e}"


@mcp.tool()
def create_presentation_from_template(filename: str, template_path: str) -> str:
    """MCP Tool: create_presentation_from_template

    Creates a new presentation based on an existing PowerPoint template (.pptx).
    This is the most reliable way to get true themes, master slides, and fonts.
    """
    try:
        if not os.path.isfile(template_path):
            return f"? ERROR: Template not found at '{template_path}'"
        abs_out = ppt_manager.create_from_template(filename, template_path)
        return f"? SUCCESS: Created from template '{template_path}' at '{abs_out}'"
    except Exception as e:
        return f"? ERROR: Failed to create from template. Details: {e}"


@mcp.tool()
def add_title_slide(filename: str, title: str, subtitle: str) -> str:
    """MCP Tool: add_title_slide

    Adds a styled cover slide with a title + subtitle.
    """
    try:
        ppt_manager.load_or_create(filename)
        slide_layout = ppt_manager._presentation.slide_layouts[0]  # Title Slide
        slide = ppt_manager._presentation.slides.add_slide(slide_layout)

        _apply_slide_background(slide)
        _add_accent_bar(slide)

        # Title
        slide.shapes.title.text = title
        _style_title_shape(slide.shapes.title)

        # Subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        tf = subtitle_shape.text_frame
        if tf.paragraphs:
            p = tf.paragraphs[0]
            p.font.size = PPTStyle.SUBTITLE_SIZE
            p.font.color.rgb = PPTStyle.MUTED

        abs_path = ppt_manager.get_absolute_path(filename)
        ppt_manager._presentation.save(abs_path)
        return f"? ADDED TITLE SLIDE: '{title}' (Saved to {abs_path})."
    except Exception as e:
        return f"? ERROR: Failed to add title slide. Details: {e}"

@mcp.tool()
def add_slide_with_title_and_bullets(filename: str, title: str, bullet_points: list[str]) -> str:
    """
    MCP Tool: add_slide_with_title_and_bullets
    Adds a standard layout slide featuring a large title and a sequential list of bullet points.
    """
    try:
        ppt_manager.load_or_create(filename)
        # Layout 1 corresponds to "Title and Content" in standard PowerPoint templates
        slide_layout = ppt_manager._presentation.slide_layouts[1]
        slide = ppt_manager._presentation.slides.add_slide(slide_layout)
        
        _apply_slide_background(slide)
        _add_accent_bar(slide)

        # 1. Format the Title Shape
        slide.shapes.title.text = title
        _style_title_shape(slide.shapes.title)
        
        # 2. Format the Main Body Content
        body_shape = slide.placeholders[1]
        body_shape.text_frame.clear()  # Drop default placeholder text
        
        # Iterate and inject cleanly formatted points
        for i, bullet in enumerate(bullet_points):
            # Take the first paragraph if index 0, else spawn a new one
            p = body_shape.text_frame.paragraphs[0] if i == 0 else body_shape.text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = PPTStyle.BODY_SIZE
            p.font.color.rgb = PPTStyle.BODY
            
        # 3. Synchronize to Disk
        abs_path = ppt_manager.get_absolute_path(filename)
        ppt_manager._presentation.save(abs_path)
        
        return f"? ADDED SLIDE: '{title}' with {len(bullet_points)} bullets (Saved to {abs_path})."
    except Exception as e:
        return f"? ERROR: Failed to add slide '{title}'. Details: {e}"

@mcp.tool()
def add_image_placeholder_slide(filename: str, title: str, image_keyword: str) -> str:
    """
    MCP Tool: add_image_placeholder_slide
    Adds a slide with a large grey rectangular placeholder indicating an expected image.
    Used when the presentation requires visual assets that the user will drop in later.
    """
    try:
        ppt_manager.load_or_create(filename)
        # Layout 5 is generally "Title Only" or purely blank
        slide_layout = ppt_manager._presentation.slide_layouts[5]
        slide = ppt_manager._presentation.slides.add_slide(slide_layout)

        _apply_slide_background(slide)
        _add_accent_bar(slide)
        
        # Safe Title Injection
        try:
            slide.shapes.title.text = title
        except Exception:
             # Fallback textbox creation if the master template layout lacks a title placeholder
             slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1)).text_frame.text = title
             
        # Draw a custom geometric rectangle dead-center
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(1.5), Inches(2.0), Inches(7.0), Inches(4.5)
        )
        
        # Format geometry with a solid light-grey finish
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(226, 232, 240)
        placeholder.line.color.rgb = RGBColor(203, 213, 225)
        
        # Imprint placeholder instruction text into the center of the rectangle
        tf = placeholder.text_frame
        p = tf.paragraphs[0]
        p.text = f"[ IMAGE PLACEHOLDER: {image_keyword} ]"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PPTStyle.MUTED
        p.alignment = PP_ALIGN.CENTER
        
        # Flush to Disk
        abs_path = ppt_manager.get_absolute_path(filename)
        ppt_manager._presentation.save(abs_path)
        return f"? ADDED IMAGE SLIDE for '{image_keyword}' (Saved to {abs_path})"
    except Exception as e:
        return f"? ERROR: Could not add image slide. Data: {e}"

@mcp.tool()
def add_two_column_content_slide(filename: str, title: str, left_bullets: list[str], right_bullets: list[str]) -> str:
    """
    MCP Tool: add_two_column_content_slide
    Adds a slide with a central title and two parallel columns of text (e.g. Pros vs Cons).
    """
    try:
        ppt_manager.load_or_create(filename)
        # Layout 3 corresponds to "Two Content"
        slide_layout = ppt_manager._presentation.slide_layouts[3]
        slide = ppt_manager._presentation.slides.add_slide(slide_layout)
        
        _apply_slide_background(slide)
        _add_accent_bar(slide)

        # Configure Title
        slide.shapes.title.text = title
        _style_title_shape(slide.shapes.title)
        
        # Locate Left and Right content boxes
        left_body = slide.placeholders[1]
        right_body = slide.placeholders[2]
        
        left_body.text_frame.clear()
        right_body.text_frame.clear()
        
        # Populate Left Column
        for i, bullet in enumerate(left_bullets):
            p = left_body.text_frame.paragraphs[0] if i == 0 else left_body.text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = PPTStyle.BODY_SIZE
            p.font.color.rgb = PPTStyle.BODY
            
        # Populate Right Column
        for i, bullet in enumerate(right_bullets):
            p = right_body.text_frame.paragraphs[0] if i == 0 else right_body.text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = PPTStyle.BODY_SIZE
            p.font.color.rgb = PPTStyle.BODY
            
        # Flush to Disk
        abs_path = ppt_manager.get_absolute_path(filename)
        ppt_manager._presentation.save(abs_path)
        return f"? ADDED TWO-COLUMN SLIDE: '{title}' (Saved to {abs_path})"
    except Exception as e:
        return f"? ERROR: Could not add two-column slide. Data: {e}"


@mcp.tool()
def add_image_slide_from_path(filename: str, title: str, image_path: str, caption: str = "") -> str:
    """MCP Tool: add_image_slide_from_path

    Adds a styled slide with a title and a locally available image.
    Use this when the user provides an image file path.
    """
    try:
        if not os.path.isfile(image_path):
            return f"? ERROR: Image file not found at '{image_path}'"

        ppt_manager.load_or_create(filename)
        slide_layout = ppt_manager._presentation.slide_layouts[5]  # Title Only
        slide = ppt_manager._presentation.slides.add_slide(slide_layout)

        _apply_slide_background(slide)
        _add_accent_bar(slide)

        # Title
        try:
            slide.shapes.title.text = title
            _style_title_shape(slide.shapes.title)
        except Exception:
            tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(9.0), Inches(0.8))
            tb.text_frame.text = title
            _style_title_shape(tb)

        # Image
        slide.shapes.add_picture(image_path, Inches(1.0), Inches(1.6), width=Inches(8.5))

        if caption:
            cap = slide.shapes.add_textbox(Inches(1.0), Inches(6.9), Inches(8.5), Inches(0.5))
            cap_tf = cap.text_frame
            cap_tf.text = caption
            p = cap_tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = PPTStyle.MUTED

        abs_path = ppt_manager.get_absolute_path(filename)
        ppt_manager._presentation.save(abs_path)
        return f"? ADDED IMAGE SLIDE: '{title}' (Saved to {abs_path})."
    except Exception as e:
        return f"? ERROR: Failed to add image slide. Details: {e}"


if __name__ == "__main__":
    # Start the MCP server using robust standard input/output streams.
    # The custom robust runner ensures empty lines/newlines over stdout don't break JSON-RPC parsing.
    from mcp_stdio_robust import run_fastmcp_with_robust_stdio

    run_fastmcp_with_robust_stdio(mcp)

